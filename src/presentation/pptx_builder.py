"""
PPTX Builder – injects analysis data into PortfolioLens_Master_Template.pptx

The template has 12 slides with {{placeholder}} text boxes.
Strategy:
  1. Open the template (preserves all formatting/layout/design)
  2. Scan every shape for {{key}} patterns
  3. Text placeholders → replace with computed values
  4. Chart placeholders → replace text box with PNG image at same position
  5. Instruction/grey-hint shapes → leave untouched

PPTX generation is blocked if qa_errors is non-empty.
"""

from __future__ import annotations

import io
import logging
import os
import tempfile

from pptx import Presentation
from pptx.util import Pt

from src.config.settings import get_settings
from src.domain.models import AnalysisOutputs
from src.utils.chart_utils import pie_chart, bar_chart
from src.utils.rtl_utils import reshape_hebrew, apply_rtl_paragraph
from src.utils.formatters import (
    fmt_ils, fmt_pct, fmt_pct_raw, fmt_duration,
)

logger = logging.getLogger(__name__)
settings = get_settings()

# Instruction shapes in the template start with these Hebrew/English strings.
# We never overwrite them – they are design-time notes for the analyst.
INSTRUCTION_PREFIXES = (
    "אסור", "יש ל", "אם אין", "אם קיימ", "אם 100%",
    "אם יש", "שקף זה", "Versioned", "להציג", "לציין",
)


class PPTXBuilder:
    """Fill the PortfolioLens Master Template with analysis data."""

    def __init__(self) -> None:
        self._template_path = settings.pptx_template_path

    # ── Public API ─────────────────────────────────────────────────────────────

    def build(self, outputs: AnalysisOutputs) -> bytes:
        if outputs.qa_errors:
            raise ValueError(
                f"Cannot generate PPTX – QA errors: " + "; ".join(outputs.qa_errors)
            )

        prs   = self._load_template()
        data  = self._build_data_map(outputs)
        charts = self._build_chart_map(outputs)

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if any(text.startswith(p) for p in INSTRUCTION_PREFIXES):
                    continue
                if "{{" not in text:
                    continue
                key = text.strip("{} \n")
                if key in charts:
                    self._replace_with_image(slide, shape, charts[key])
                elif key in data:
                    self._replace_text(shape, data[key])

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    # ── Template loading ───────────────────────────────────────────────────────

    def _load_template(self) -> Presentation:
        if self._template_path.exists():
            logger.info(f"Loading template: {self._template_path}")
            return Presentation(str(self._template_path))
        logger.warning("Template not found – using blank fallback")
        return Presentation()

    # ── Data map ───────────────────────────────────────────────────────────────

    def _build_data_map(self, o: AnalysisOutputs) -> dict:
        prefs = o.preferences

        eq_row  = next((r for r in o.asset_allocation if r.asset_class == "equity"), None)
        bnd_row = next((r for r in o.asset_allocation if r.asset_class == "bond"),   None)
        csh_row = next((r for r in o.asset_allocation if r.asset_class == "cash"),   None)
        eq_pct  = fmt_pct(eq_row.weight)  if eq_row  else "—"
        bnd_pct = fmt_pct(bnd_row.weight) if bnd_row else "—"
        csh_pct = fmt_pct(csh_row.weight) if csh_row else "—"

        # Asset allocation table
        aa_lines = ["סוג נכס | שווי (ILS) | משקל"]
        for r in o.asset_allocation:
            aa_lines.append(f"{r.asset_class} | {fmt_ils(r.market_value_ils)} | {fmt_pct(r.weight)}")

        # Equity geography
        total_eq = sum(r.market_value_ils for r in o.equity_geography) or 1
        il_row = next((r for r in o.equity_geography if "Israel" in r.region or "ישראל" in r.region), None)
        foreign_val = sum(r.market_value_ils for r in o.equity_geography if "Israel" not in r.region and "ישראל" not in r.region)
        dom_text = (
            f"ישראל: {fmt_ils(il_row.market_value_ils)} ({fmt_pct(il_row.weight_in_equities)})\n"
            f"חו\"ל: {fmt_ils(foreign_val)} ({fmt_pct(foreign_val/total_eq)})"
        ) if il_row else f"חו\"ל: {fmt_ils(foreign_val)} ({fmt_pct(foreign_val/total_eq)})"

        geo_lines = [f"{r.region}: {fmt_ils(r.market_value_ils)} | {fmt_pct(r.weight_in_equities)}" + (" [E]" if r.is_estimated else "") for r in o.equity_geography]

        # US exposure
        us = o.us_exposure
        us_cons  = fmt_pct(us.conservative_us_weight) if us else "—"
        us_broad = fmt_pct(us.broad_us_weight)        if us else "—"

        # Sector
        sector_lines = ["ענף | שווי (ILS) | % ממניות"]
        for r in o.sector_allocation:
            sector_lines.append(f"{r.sector} | {fmt_ils(r.market_value_ils)} | {fmt_pct(r.weight_in_equities)}" + (" [E]" if r.is_estimated else ""))

        # Bond
        bond_lines = ["סוג | שווי (ILS) | % מאגח | % מתיק"]
        for r in o.bond_breakdown:
            bond_lines.append(f"{r.linkage_type} | {fmt_ils(r.market_value_ils)} | {fmt_pct(r.weight_in_bonds)} | {fmt_pct(r.weight_in_portfolio)}")

        # Duration
        dur_lines = ["נייר | שווי | מח\"מ | מקור"]
        for r in o.duration_table:
            dur_lines.append(f"{r.name[:25]} | {fmt_ils(r.market_value_ils)} | {fmt_duration(r.duration)}" + (" [E]" if r.is_estimated else "") + f" | {r.duration_source[:20]}")

        # Fund costs
        cost_lines = ["קרן | משקל | דמי ניהול | עלות שנתית"]
        for r in o.fund_cost_table:
            cost_lines.append(f"{r.name[:25]} | {fmt_pct(r.weight_in_portfolio)} | {fmt_pct_raw(r.fee_percent)}" + (" [E]" if r.is_estimated else "") + f" | {fmt_ils(r.annual_cost_ils) if r.annual_cost_ils else '—'}")

        pm_fee_str = (
            fmt_pct_raw(prefs.portfolio_manager_fee_percent)
            + (" [ASSUMPTION]" if prefs.manager_fee_is_assumption else "")
            if prefs and prefs.portfolio_manager_fee_percent else "לא הוזן"
        )
        total_cost_str = (
            fmt_pct_raw(o.total_cost_percent) + (" [ASSUMPTION]" if o.total_cost_is_assumption else "")
            if o.total_cost_percent else "לא חושב (חסר דמי מנהל)"
        )

        # FX
        fx_lines = ["מטבע | שווי | % | גידור"]
        for r in o.fx_exposure:
            fx_lines.append(f"{r.currency} | {fmt_ils(r.market_value_ils)} | {fmt_pct(r.weight)} | {r.hedging_note}")
        fx_insight = (
            f"ILS: {fmt_pct(next((r.weight for r in o.fx_exposure if r.currency == 'ILS'), 0))} | "
            f"FX: {fmt_pct(sum(r.weight for r in o.fx_exposure if r.currency != 'ILS'))}"
            if o.fx_exposure else "—"
        )

        # Top holdings
        top_lines = ["# | שם | שווי | משקל"]
        for r in o.top_holdings:
            top_lines.append(f"{r.rank} | {r.name[:30]} | {fmt_ils(r.market_value_ils)} | {fmt_pct(r.weight)}")

        # Assumptions
        assume_lines = ["שדה | נייר | ערך | סיבה"]
        for a in o.assumptions[:15]:
            assume_lines.append(f"{a.field} | {a.holding_name[:20]} | {a.assumed_value} | {a.reason[:30]}")

        # QA
        qa_parts = (
            [f"❌ {e}" for e in o.qa_errors] +
            [f"⚠️ {w}" for w in o.qa_warnings] or
            ["✅ כל בדיקות ה-QA עברו בהצלחה"]
        )

        # Summary bullets
        bullets = [
            f"• שווי תיק: {fmt_ils(o.total_portfolio_value_ils)}",
            f"• מניות: {eq_pct} | אגח: {bnd_pct} | מזומן: {csh_pct}",
        ]
        if o.conservative_weighted_duration:
            bullets.append(f"• מח\"מ שמרן: {fmt_duration(o.conservative_weighted_duration)}")
        if o.effective_fund_cost_on_total_portfolio:
            bullets.append(f"• עלות קרנות: {fmt_pct_raw(o.effective_fund_cost_on_total_portfolio)}")
        if o.top5_concentration:
            bullets.append(f"• Top-5 ריכוזיות: {fmt_pct(o.top5_concentration)}")

        return {
            # Cover
            "portfolio_value":  fmt_ils(o.total_portfolio_value_ils),
            "analysis_date":    prefs.report_date if prefs else "",
            "data_date":        prefs.report_date if prefs else "",
            "client_name":      prefs.client_name if prefs else "",
            # Executive summary
            "equity_percent":   eq_pct,
            "bond_percent":     bnd_pct,
            "cash_percent":     csh_pct,
            "summary_bullets":  "\n".join(bullets),
            "summary_visual":   "",
            # Asset allocation
            "table_asset_allocation":       "\n".join(aa_lines),
            "chart_asset_allocation_donut": "",
            "insight_asset_allocation":     f"סה\"כ {len(o.asset_allocation)} קטגוריות נכסים",
            "method_asset_allocation":      "מזומן " + ("נכלל" if prefs and prefs.include_cash_in_allocation else "לא נכלל") + " בחישוב",
            # Equity geography
            "table_equity_domestic_foreign":    dom_text,
            "table_equity_geography_breakdown": "\n".join(geo_lines),
            "chart_equity_geography":           "",
            "insight_equity_geography":         geo_lines[0] if geo_lines else "—",
            # US Exposure
            "us_exposure_conservative":     us_cons,
            "us_exposure_broad":            us_broad,
            "us_exposure_methodology":      us.methodology_note if us else "—",
            "chart_us_exposure_comparison": "",
            "insight_us_exposure":          f"Conservative: {us_cons} | Broad: {us_broad}",
            # Sector
            "table_sector_allocation":      "\n".join(sector_lines),
            "chart_sector_treemap":         "",
            "insight_sector_allocation":    (f"מוביל: {o.sector_allocation[0].sector} ({fmt_pct(o.sector_allocation[0].weight_in_equities)})" if o.sector_allocation else "—"),
            "footnote_sector_methodology":  "breakdowns רשמיים + משקלות מדד כ-proxy לקרנות",
            # Bond
            "table_bond_breakdown":         "\n".join(bond_lines),
            "chart_bond_breakdown":         "",
            "insight_bond_breakdown":       (f"סה\"כ אגח: {fmt_ils(sum(r.market_value_ils for r in o.bond_breakdown))}" if o.bond_breakdown else "—"),
            "method_bond_breakdown":        "סיווג לפי שם + מטבע + ISIN",
            # Duration
            "table_duration":                    "\n".join(dur_lines),
            "chart_duration_bar":                "",
            "weighted_duration_conservative":    fmt_duration(o.conservative_weighted_duration),
            "weighted_duration_extended":        fmt_duration(o.extended_weighted_duration),
            "method_duration":                   "שמרן: מקורות רשמיים בלבד. מורחב: כולל הערכות.",
            # Cost
            "table_fund_costs":                              "\n".join(cost_lines),
            "weighted_fund_cost":                            fmt_pct_raw(o.weighted_fund_cost_on_funds),
            "effective_fund_cost_on_total_portfolio":        fmt_pct_raw(o.effective_fund_cost_on_total_portfolio),
            "portfolio_manager_fee":                         pm_fee_str,
            "total_cost_percent":                            total_cost_str,
            "cost_notes":                                    "; ".join(o.methodology_notes[:2]),
            # FX
            "table_fx_exposure":    "\n".join(fx_lines),
            "chart_fx_exposure":    "",
            "insight_fx_exposure":  fx_insight,
            "method_fx_exposure":   "גידור לפי מדיניות קרן / הצהרת מנהל",
            # Concentration
            "table_top_holdings":       "\n".join(top_lines),
            "top_5_weight":             fmt_pct(o.top5_concentration),
            "top_10_weight":            fmt_pct(o.top10_concentration),
            "single_name_risk_note":    "\n".join(o.concentration_warnings) if o.concentration_warnings else "אין ריכוזיות חריגה",
            "chart_concentration":      "",
            # Appendix
            "table_assumptions":    "\n".join(assume_lines) if o.assumptions else "אין הנחות — כל הנתונים ממקורות מאומתים",
            "table_sources":        "\n".join(o.source_urls[:10]) if o.source_urls else "Mock research provider",
            "qa_checklist":         "\n".join(qa_parts),
        }

    # ── Chart map ──────────────────────────────────────────────────────────────

    def _build_chart_map(self, o: AnalysisOutputs) -> dict:
        c: dict[str, bytes] = {}

        if o.asset_allocation:
            img = pie_chart([r.asset_class for r in o.asset_allocation], [r.weight*100 for r in o.asset_allocation], figsize=(4,3))
            c["chart_asset_allocation_donut"] = img
            c["summary_visual"] = img

        if o.equity_geography:
            c["chart_equity_geography"] = pie_chart([r.region for r in o.equity_geography], [r.weight_in_equities*100 for r in o.equity_geography], figsize=(4,3))

        if o.us_exposure:
            c["chart_us_exposure_comparison"] = bar_chart(["Conservative","Broad"], [o.us_exposure.conservative_us_weight*100, o.us_exposure.broad_us_weight*100], horizontal=False, figsize=(3,3), value_format="{:.1f}%")

        if o.sector_allocation:
            c["chart_sector_treemap"] = bar_chart([r.sector for r in o.sector_allocation[:8]], [r.weight_in_equities*100 for r in o.sector_allocation[:8]], horizontal=True, figsize=(4,3))

        if o.bond_breakdown:
            c["chart_bond_breakdown"] = pie_chart([r.linkage_type for r in o.bond_breakdown], [r.weight_in_bonds*100 for r in o.bond_breakdown], figsize=(4,3))

        dur_data = [(r.name[:20], r.duration) for r in o.duration_table if r.duration is not None]
        if dur_data:
            c["chart_duration_bar"] = bar_chart([d[0] for d in dur_data], [d[1] for d in dur_data], horizontal=True, figsize=(4,3), value_format="{:.1f}y")

        if o.fx_exposure:
            c["chart_fx_exposure"] = pie_chart([r.currency for r in o.fx_exposure], [r.weight*100 for r in o.fx_exposure], figsize=(4,3))

        if o.top_holdings:
            c["chart_concentration"] = bar_chart([r.name[:20] for r in o.top_holdings[:10]], [r.weight*100 for r in o.top_holdings[:10]], horizontal=True, figsize=(4,3))

        return c

    # ── Shape helpers ──────────────────────────────────────────────────────────

    @staticmethod
    def _replace_text(shape, value: str) -> None:
        tf = shape.text_frame
        try:
            p0  = tf.paragraphs[0]
            r0  = p0.runs[0] if p0.runs else None
            sz  = r0.font.size  if r0 else None
            bld = r0.font.bold  if r0 else None
            try:
                clr = r0.font.color.rgb if (r0 and r0.font.color and r0.font.color.type) else None
            except Exception:
                clr = None
        except Exception:
            sz = bld = clr = None

        tf.clear()
        para = tf.paragraphs[0]
        run  = para.add_run()
        run.text = reshape_hebrew(value) if value else "—"
        if sz:   run.font.size  = sz
        if bld is not None: run.font.bold = bld
        if clr:  run.font.color.rgb = clr
        apply_rtl_paragraph(para)

    @staticmethod
    def _replace_with_image(slide, shape, img_bytes: bytes) -> None:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        try:
            shape.text_frame.clear()
        except Exception:
            pass
        tmp_fd, tmp_path = tempfile.mkstemp(suffix=".png")
        try:
            with os.fdopen(tmp_fd, "wb") as f:
                f.write(img_bytes)
            slide.shapes.add_picture(tmp_path, left, top, width=width, height=height)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

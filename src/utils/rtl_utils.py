"""
RTL / Hebrew utilities for correct text rendering in PowerPoint.

python-pptx does not natively support RTL paragraphs.
These helpers apply the necessary XML attributes and BiDi reshaping
so that Hebrew text is stored and rendered correctly.

Usage
-----
from src.utils.rtl_utils import apply_rtl_paragraph, reshape_hebrew

shape = slide.shapes.add_textbox(...)
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = reshape_hebrew("שלום עולם")
apply_rtl_paragraph(p)
"""

from __future__ import annotations

import re
from lxml import etree


def reshape_hebrew(text: str) -> str:
    """
    Apply arabic-reshaper + python-bidi to ensure RTL characters are
    stored in the correct visual order for embedding in PPTX XML.

    Falls back gracefully if libraries are not installed.
    """
    try:
        import arabic_reshaper
        from bidi.algorithm import get_display
        reshaped = arabic_reshaper.reshape(text)
        return get_display(reshaped)
    except ImportError:
        # Libraries not installed – return as-is (text may render reversed in old viewers)
        return text


def apply_rtl_paragraph(paragraph) -> None:
    """
    Set the paragraph-level RTL XML attribute on a python-pptx Paragraph object.
    This makes PowerPoint render the paragraph right-to-left.
    """
    pPr = paragraph._p.get_or_add_pPr()
    # Set rtl="1" attribute in the DrawingML namespace
    pPr.set("rtl", "1")


def apply_rtl_run(run) -> None:
    """
    Set RTL attribute on a Run (character-level BiDi override).
    """
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "he-IL")


def make_rtl_paragraph(text_frame, text: str, font_size_pt: int = 12, bold: bool = False) -> None:
    """
    Convenience: clear text_frame, add one RTL paragraph with given text.
    """
    from pptx.util import Pt
    tf = text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = reshape_hebrew(text)
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    apply_rtl_paragraph(p)
    apply_rtl_run(run)


def is_rtl_text(text: str) -> bool:
    """Detect if a string contains RTL characters (Hebrew, Arabic)."""
    RTL_PATTERN = re.compile(r"[\u0590-\u05FF\u0600-\u06FF]")
    return bool(RTL_PATTERN.search(text))


def safe_fill_placeholder(placeholder, text: str, rtl: bool = True) -> None:
    """
    Fill a PPTX placeholder with text, applying RTL if needed.
    Handles cases where placeholder has no text frame gracefully.
    """
    try:
        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = reshape_hebrew(text) if rtl else text
        if rtl or is_rtl_text(text):
            apply_rtl_paragraph(p)
            apply_rtl_run(run)
    except Exception:
        # Some placeholders don't have text frames (image placeholders, etc.)
        try:
            placeholder.text = text
        except Exception:
            pass  # Cannot fill this placeholder type
